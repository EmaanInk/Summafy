
# Import libraries for reading different file types
import pymupdf as fitz
from docx import Document
from pptx import Presentation


# Function to read text from PDF files
def read_pdf(file_path):

    # Empty string to store extracted text
    text = ""

    # Open the PDF document
    doc = fitz.open(file_path)

    # Loop through all pages in PDF
    for pages in doc:

        # Extract text from each page and add it to text variable
        text += pages.get_text()

    # Return complete extracted text
    return text


# Function to read text from TXT files
def read_txt(file_path):

    # Open file in read mode with UTF-8 encoding
    with open(file_path, 'r', encoding='utf-8') as file:

        # Read all file content
        data = file.read()

    # Return text data
    return data


# Function to read text from DOCX files
def read_docx(file_path):

    # Empty string to store extracted text
    text = ""

    # Open Word document
    doc = Document(file_path)

    # Loop through all paragraphs
    for paragraph in doc.paragraphs:

        # Add paragraph text to text variable
        text += paragraph.text

    # Return extracted text
    return text


# Function to read text from PPTX files
def read_pptx(file_path):

    # Empty string to store extracted text
    text = ""

    # Open PowerPoint presentation
    doc = Presentation(file_path)

    # Loop through all slides
    for slide in doc.slides:

        # Loop through all shapes in slide
        for shapes in slide.shapes:

            # Check if shape contains text
            if shapes.has_text_frame:

                # Add shape text to text variable
                text += shapes.text

    # Return extracted text
    return text


# Main function to detect file type and call correct reader function
def read_files(file_path):

    # Check if file is PDF
    if file_path.endswith(".pdf"):
        return read_pdf(file_path)

    # Check if file is DOCX
    elif file_path.endswith(".docx"):
        return read_docx(file_path)

    # Check if file is TXT
    elif file_path.endswith(".txt"):
        return read_txt(file_path)

    # Check if file is PPTX
    elif file_path.endswith(".pptx"):
        return read_pptx(file_path)

    # If file type is unsupported
    else:
        return "Unsupported file Type"

